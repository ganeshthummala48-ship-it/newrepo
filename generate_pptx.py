from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Helper function to add slide with title and content
    def add_slide(title, content):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]

        title_shape.text = title
        tf = body_shape.text_frame
        tf.text = content

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Farmer AI"
    subtitle.text = "Empowering the Future of Agriculture\nAn AI-Driven Integrated Mobile Solution"

    # Slide 2: Problem Statement
    add_slide("Problem Statement", 
              "• Farmers face significant language barriers when accessing advanced agricultural info.\n"
              "• Delay in crop disease detection leads to major harvest losses.\n"
              "• Lack of direct market access and fair pricing for services and machinery.\n"
              "• Difficulty in navigating government schemes and subsidies.")

    # Slide 3: Existing System
    add_slide("Existing System", 
              "• Manual diagnosis of plant diseases based on visual experience.\n"
              "• Information seeking through local shops or government offices (often slow).\n"
              "• Fragmented communication for hiring machinery or labour.\n"
              "• Paper-based or complex web portals for government schemes.")

    # Slide 4: Drawbacks of Existing System
    add_slide("Drawbacks of Existing System", 
              "• High error rates in manual disease diagnosis.\n"
              "• Late response to pest outbreaks causing irreversible crop damage.\n"
              "• Inefficient resource utilization and lack of transparency in pricing.\n"
              "• Exclusion of farmers who are not tech-savvy or fluent in English.")

    # Slide 5: Proposed System
    add_slide("Proposed System: Farmer AI", 
              "• Multi-lingual mobile application (Telugu, Hindi, English).\n"
              "• AI-powered real-time crop disease detection via image recognition.\n"
              "• Intelligent crop and scheme recommendation engine.\n"
              "• Direct marketplace for machinery, labour, and fertilizers with AI negotiation support.")

    # Slide 6: Technologies Used
    add_slide("Technologies Used", 
              "• Frontend: Flutter (Cross-platform mobile framework).\n"
              "• Backend: FastAPI (High-performance Python API).\n"
              "• AI Models: TensorFlow (Disease detection), Cohere API (LLM for recommendations).\n"
              "• Database: SQLite (Lightweight, efficient data storage).")

    # Slide 7: Flow Chart
    add_slide("Flow Chart: User Journey", 
              "1. User Registration/Login (localized profile selection).\n"
              "2. Image Upload -> AI Disease Detection -> Real-time treatment recommendation.\n"
              "3. Input Conditions -> AI Crop/Scheme recommendation.\n"
              "4. Service Listing/Search -> Integrated AI Negotiation Tool.")

    # Slide 8: Conclusion
    add_slide("Conclusion", 
              "• Farmer AI bridging the gap between advanced technology and rural agriculture.\n"
              "• Increased crop yields through early disease detection and smart planning.\n"
              "• Empowering farmers with localized info and fair market access.\n"
              "• Future Scope: IoT integration and real-time weather analytics.")

    # Save the presentation
    prs.save('Farmer_AI_Presentation.pptx')
    print("Presentation created successfully: Farmer_AI_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
